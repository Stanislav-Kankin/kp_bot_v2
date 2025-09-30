from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
import os
from datetime import datetime
import re
import logging

class PPTService:
    def __init__(self):
        self.templates_path = "templates/"
        # Основной шрифт для всего документа
        self.primary_font = 'Montserrat'
        # Запасные шрифты на случай отсутствия Montserrat
        self.fallback_fonts = ['Arial', 'Helvetica', 'Times New Roman', 'Calibri']
        self.logger = logging.getLogger('PPTService')
    
    def _ensure_consistent_fonts(self, prs):
        """Устанавливает единый шрифт Montserrat для всего текста в презентации"""
        for slide in prs.slides:
            for shape in slide.shapes:
                self._set_shape_font(shape)
    
    def _set_shape_font(self, shape):
        """Устанавливает шрифт Montserrat для фигуры и всего её содержимого"""
        if not shape.has_text_frame:
            return
            
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Сохраняем оригинальные свойства форматирования
                original_size = run.font.size
                original_bold = run.font.bold
                original_italic = run.font.italic
                original_color = run.font.color.rgb if run.font.color and run.font.color.rgb else None
                original_underline = run.font.underline
                
                # Устанавливаем шрифт Montserrat
                run.font.name = self.primary_font
                
                # Восстанавливаем форматирование
                if original_size:
                    run.font.size = original_size
                run.font.bold = original_bold
                run.font.italic = original_italic
                run.font.underline = original_underline
                if original_color:
                    run.font.color.rgb = original_color
    
    def _set_table_font(self, table):
        """Устанавливает шрифт Montserrat для всей таблицы"""
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    text_frame = cell.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Сохраняем оригинальные свойства
                            original_size = run.font.size
                            original_bold = run.font.bold
                            original_italic = run.font.italic
                            original_color = run.font.color.rgb if run.font.color and run.font.color.rgb else None
                            
                            # Устанавливаем шрифт Montserrat
                            run.font.name = self.primary_font
                            
                            # Восстанавливаем форматирование
                            if original_size:
                                run.font.size = original_size
                            run.font.bold = original_bold
                            run.font.italic = original_italic
                            if original_color:
                                run.font.color.rgb = original_color
    
    def create_kp_presentation(self, template_type: str, data: dict):
        """Создает КП на основе шаблона с единым шрифтом Montserrat"""
        try:
            # Выбираем шаблон
            if template_type == "long":
                template_file = "kedo_long.pptx"
            else:
                template_file = "kedo_short.pptx"

            template_path = os.path.join(self.templates_path, template_file)

            if not os.path.exists(template_path):
                raise FileNotFoundError(f"Шаблон {template_file} не найден")

            # Загружаем шаблон
            prs = Presentation(template_path)
            
            # Устанавливаем единый шрифт Montserrat для всей презентации
            self._ensure_consistent_fonts(prs)

            # Заменяем название компании
            self._replace_company_name(prs, data['company_name'])

            # Сначала обновляем таблицу, чтобы рассчитать правильную сумму
            total_price = self._update_third_table(prs, data)

            # Затем обновляем текст с стоимостью над таблицей
            self._update_price_text(prs, total_price)

            # Сохраняем результат
            output_filename = f"КП_{data['company_name']}_{datetime.now().strftime('%d%m%Y_%H%M')}.pptx"
            output_path = os.path.join(
                self.templates_path, "output", output_filename)

            # Создаем папку output если ее нет
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            prs.save(output_path)
            return output_path

        except Exception as e:
            print(f"Ошибка при создании презентации: {e}")
            return None

    def _replace_company_name(self, prs, company_name):
        """Заменяет название компании на первом слайде с сохранением форматирования"""
        try:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if hasattr(shape, "text") and "Название" in shape.text:
                    self._safe_text_replace(shape, "Название", company_name)
                    return
        except Exception as e:
            print(f"Ошибка при замене названия компании: {e}")

    def _safe_text_replace(self, shape, old_text, new_text):
        """Безопасная замена текста с сохранением форматирования и установкой Montserrat"""
        if not shape.has_text_frame:
            return

        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    # Сохраняем свойства шрифта
                    original_font = run.font
                    original_size = original_font.size
                    original_bold = original_font.bold
                    original_italic = original_font.italic
                    original_color = original_font.color.rgb if original_font.color and original_font.color.rgb else None
                    original_underline = original_font.underline

                    # Заменяем текст
                    run.text = run.text.replace(old_text, new_text)

                    # Устанавливаем шрифт Montserrat
                    run.font.name = self.primary_font
                    
                    # Восстанавливаем форматирование
                    if original_size:
                        run.font.size = original_size
                    run.font.bold = original_bold
                    run.font.italic = original_italic
                    run.font.underline = original_underline
                    if original_color:
                        run.font.color.rgb = original_color

    def _update_price_text(self, prs, total_price):
        """Обновляет текст с стоимостью над таблицей с сохранением форматирования"""
        try:
            formatted_price = f"{total_price:,} ₽".replace(',', ' ')
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and "Стоимость HRlink на 12 месяцев" in shape.text:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    # Ищем и заменяем только цифры с ₽
                                    if re.search(r'\d[\d\s]*₽', run.text):
                                        original_font = run.font
                                        original_size = original_font.size
                                        original_bold = original_font.bold
                                        original_italic = original_font.italic
                                        original_color = original_font.color.rgb if original_font.color and original_font.color.rgb else None
                                        original_underline = original_font.underline

                                        # Заменяем всё число с ₽
                                        run.text = re.sub(
                                            r'\d[\d\s]*₽',
                                            formatted_price,
                                            run.text
                                        )

                                        # Устанавливаем шрифт Montserrat
                                        run.font.name = self.primary_font
                                        
                                        # Восстанавливаем форматирование
                                        if original_size:
                                            run.font.size = original_size
                                        run.font.bold = original_bold
                                        run.font.italic = original_italic
                                        run.font.underline = original_underline
                                        if original_color:
                                            run.font.color.rgb = original_color
                            return
        except Exception as e:
            print(f"Ошибка при обновлении цены: {e}")

    def _update_third_table(self, prs, data):
        """Обновляет третью таблицу в презентации и возвращает итоговую сумму"""
        try:
            table_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        table_count += 1
                        # Третья таблица (таблица с расчетами)
                        if table_count == 3:
                            table = shape.table
                            # Устанавливаем шрифт для таблицы
                            self._set_table_font(table)
                            return self._fill_calculation_table(table, data)
            return 0
        except Exception as e:
            print(f"Ошибка при обновлении таблицы: {e}")
            return 0

    def _fill_calculation_table(self, table, data):
        prices = {
            'base': 15000,
            'hr': 15000,
            'employee': 1000,
            'on_premise': 600000
        }

        # Рассчитываем суммы
        base_total = prices['base']
        hr_total = prices['hr'] * data['hr_licenses']
        employee_total = prices['employee'] * data['employee_licenses']
        on_premise_total = prices['on_premise'] if data['on_premises'] == 'Да' else 0

        # Итоговая сумма
        total_price = base_total + hr_total + employee_total + on_premise_total

        # Обновляем таблицу
        for i, row in enumerate(table.rows):
            row_text = ' '.join([cell.text.strip() for cell in row.cells]).lower()

            if 'базовая' in row_text:
                if len(row.cells) >= 3:
                    self._safe_cell_replace(row.cells[2], row.cells[2].text, "1 шт")
                if len(row.cells) >= 5:
                    self._safe_cell_replace(row.cells[4], row.cells[4].text, f"{base_total:,} ₽".replace(',', ' '))

            elif any(word in row_text for word in ['кадровик', 'кадровика']):
                if len(row.cells) >= 3:
                    self._safe_cell_replace(row.cells[2], row.cells[2].text, f"{data['hr_licenses']} шт")
                if len(row.cells) >= 5:
                    self._safe_cell_replace(row.cells[4], row.cells[4].text, f"{hr_total:,} ₽".replace(',', ' '))

            elif any(word in row_text for word in ['сотрудник', 'сотрудника']):
                if len(row.cells) >= 3:
                    self._safe_cell_replace(row.cells[2], row.cells[2].text, f"{data['employee_licenses']} шт")
                if len(row.cells) >= 5:
                    self._safe_cell_replace(row.cells[4], row.cells[4].text, f"{employee_total:,} ₽".replace(',', ' '))

            elif 'on-premise' in row_text:
                quantity = 1 if data['on_premises'] == 'Да' else 0
                if len(row.cells) >= 3:
                    self._safe_cell_replace(row.cells[2], row.cells[2].text, f"{quantity} шт")
                if len(row.cells) >= 5:
                    self._safe_cell_replace(row.cells[4], row.cells[4].text, f"{on_premise_total:,} ₽".replace(',', ' '))

            elif any(word in row_text for word in ['итог', 'итого', 'итоговая']):
                if len(row.cells) >= 5:
                    self._safe_cell_replace(row.cells[4], row.cells[4].text, f"{total_price:,} ₽".replace(',', ' '))

        return total_price

    def _safe_cell_replace(self, cell, old_text, new_text):
        """Безопасная замена текста в ячейке с сохранением форматирования"""
        if not cell.text_frame:
            return

        text_frame = cell.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text or run.text.strip():
                    # Сохраняем свойства шрифта
                    original_font = run.font
                    original_size = original_font.size
                    original_bold = original_font.bold
                    original_italic = original_font.italic
                    original_color = original_font.color.rgb if original_font.color and original_font.color.rgb else None
                    original_underline = original_font.underline

                    # Заменяем текст
                    run.text = new_text

                    # Устанавливаем шрифт Montserrat
                    run.font.name = self.primary_font
                    
                    # Восстанавливаем форматирование
                    if original_size:
                        run.font.size = original_size
                    run.font.bold = original_bold
                    run.font.italic = original_italic
                    run.font.underline = original_underline
                    if original_color:
                        run.font.color.rgb = original_color
                    return

    def convert_to_pdf(self, pptx_path: str):
        """Конвертирует PPTX в PDF с упрощенными параметрами"""
        try:
            import subprocess

            # Получаем директорию исходного файла презентации
            output_dir = os.path.dirname(pptx_path)

            # Формируем путь для PDF-файла
            pdf_filename = os.path.basename(pptx_path).replace('.pptx', '.pdf')
            pdf_path = os.path.join(output_dir, pdf_filename)

            print(f"🔄 Начинаю конвертацию: {pptx_path} -> {pdf_path}")

            # Простая команда без сложных параметров
            command = [
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', output_dir,
                pptx_path
            ]
            
            print(f"🔧 Выполняю команду: {' '.join(command)}")
            
            result = subprocess.run(
                command,
                check=True,
                capture_output=True,
                text=True,
                timeout=60
            )
            
            print(f"✅ Конвертация завершена успешно")
            print(f"📋 STDOUT: {result.stdout}")
            
            if result.stderr:
                print(f"⚠️  STDERR: {result.stderr}")

            # Проверяем, что файл создался
            if os.path.exists(pdf_path):
                print(f"📄 PDF создан: {pdf_path}")
                return pdf_path
            else:
                print(f"❌ PDF файл не найден по пути: {pdf_path}")
                return None

        except subprocess.CalledProcessError as e:
            print(f"❌ Ошибка выполнения команды: {e}")
            print(f"🔧 STDOUT: {e.stdout}")
            print(f"🔧 STDERR: {e.stderr}")
            return None
        except subprocess.TimeoutExpired:
            print("⏰ Таймаут конвертации в PDF")
            return None
        except Exception as e:
            print(f"❌ Неожиданная ошибка при конвертации в PDF: {e}")
            return None


ppt_service = PPTService()